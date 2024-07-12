import os
import PyPDF2
import fitz  # PyMuPDF
from docx import Document as DocxDocument
import xlrd  # For .xls files
from openpyxl import load_workbook
import pypandoc
from langchain.docstore.document import Document
from pptx import Presentation

def read_txt(file_path):
    with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
        return file.read()

def read_pdf(file_path):
    text = ""
    try:
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page_num in range(len(pdf_reader.pages)):
                text += pdf_reader.pages[page_num].extract_text()
    except:
        pdf_document = fitz.open(file_path)
        for page_num in range(pdf_document.page_count):
            page = pdf_document.load_page(page_num)
            text += page.get_text()
    return text

def read_docx(file_path):
    doc = DocxDocument(file_path)
    return "\n".join([paragraph.text for paragraph in doc.paragraphs])

def read_doc(file_path):
    return pypandoc.convert_file(file_path, 'plain')

def read_xls(file_path):
    workbook = xlrd.open_workbook(file_path)
    sheet = workbook.sheet_by_index(0)
    data = []
    for row_num in range(sheet.nrows):
        data.append("\t".join(map(str, sheet.row_values(row_num))))
    return "\n".join(data)

def read_xlsx(file_path):
    workbook = load_workbook(file_path)
    sheet = workbook.active
    data = []
    for row in sheet.iter_rows(values_only=True):
        data.append("\t".join(map(str, row)))
    return "\n".join(data)

def read_ppt(file_path):
    presentation = Presentation(file_path)
    text = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

def load_files(directory):
    documents = []
    unread_documents = []
    for filename in os.listdir(directory):
        file_path = os.path.join(directory, filename)
        try:
            if filename.endswith('.txt'):
                content = read_txt(file_path)
            elif filename.endswith('.pdf'):
                content = read_pdf(file_path)
            elif filename.endswith('.docx'):
                content = read_docx(file_path)
            elif filename.endswith('.doc'):
                content = read_doc(file_path)
            elif filename.endswith('.xls'):
                content = read_xls(file_path)
            elif filename.endswith('.xlsx'):
                content = read_xlsx(file_path)
            elif filename.endswith('.ppt') or filename.endswith('.pptx'):
                content = read_ppt(file_path)
            else:
                print(f"Unsupported file format: {filename}")
                unread_documents.append(filename)
                continue

            doc = Document(page_content=content, metadata={"filename": filename})
            documents.append(doc)
        except Exception as e:
            print(f"Failed to process {filename}: {e}")
            unread_documents.append(filename)
    return documents, unread_documents

# Example usage
directory = 'path_to_your_directory'
documents, unread_documents = load_files(directory)

for doc in documents:
    print(f"Filename: {doc.metadata['filename']}")
    print(f"Content: {doc.page_content[:200]}")  # Print first 200 characters of content
    print("\n")

if unread_documents:
    print("Unread Documents:")
    for filename in unread_documents:
        print(filename)

# Optionally, save the unread documents to a log file
with open('unread_documents.log', 'w') as log_file:
    for filename in unread_documents:
        log_file.write(f"{filename}\n")
