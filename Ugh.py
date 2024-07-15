import os
import PyPDF2
import fitz  # PyMuPDF
from docx import Document as DocxDocument
import xlrd  # For .xls files
from openpyxl import load_workbook
from pptx import Presentation
import mammoth
import subprocess
from langchain.docstore.document import Document
from vertexaiembeddings import VertexAIEmbeddings
from langchain_community.vectorstore import FAISS
import pickle

def read_txt(file_path):
    with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
        return file.read()

def read_pdf(file_path):
    text = ""
    try:
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page_num in range(len(pdf_reader.pages)):
                text += pdf_reader.pages[page_num].extract_text()
    except:
        pdf_document = fitz.open(file_path)
        for page_num in range(pdf_document.page_count):
            page = pdf_document.load_page(page_num)
            text += page.get_text()
    return text

def read_docx(file_path):
    doc = DocxDocument(file_path)
    return "\n".join([paragraph.text for paragraph in doc.paragraphs])

def read_docm(file_path):
    with open(file_path, 'rb') as docm_file:
        result = mammoth.extract_raw_text(docm_file)
        return result.value

def read_xls(file_path):
    workbook = xlrd.open_workbook(file_path)
    sheet = workbook.sheet_by_index(0)
    data = []
    for row_num in range(sheet.nrows):
        data.append("\t".join(map(str, sheet.row_values(row_num))))
    return "\n".join(data)

def read_xlsx(file_path):
    workbook = load_workbook(file_path)
    sheet = workbook.active
    data = []
    for row in sheet.iter_rows(values_only=True):
        data.append("\t".join(map(str, row)))
    return "\n".join(data)

def read_ppt(file_path):
    presentation = Presentation(file_path)
    text = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

def read_doc(file_path):
    result = subprocess.run(['antiword', file_path], capture_output=True, text=True)
    if result.returncode != 0:
        raise Exception(f"Error reading {file_path}: {result.stderr}")
    return result.stdout

def load_files(directory):
    documents = []
    unread_documents = []
    for filename in os.listdir(directory):
        file_path = os.path.join(directory, filename)
        try:
            if filename.endswith('.txt'):
                content = read_txt(file_path)
            elif filename.endswith('.pdf'):
                content = read_pdf(file_path)
            elif filename.endswith('.docx'):
                content = read_docx(file_path)
            elif filename.endswith('.doc'):
                content = read_doc(file_path)
            elif filename.endswith('.docm'):
                content = read_docm(file_path)
            elif filename.endswith('.xls'):
                content = read_xls(file_path)
            elif filename.endswith('.xlsx'):
                content = read_xlsx(file_path)
            elif filename.endswith('.ppt') or filename.endswith('.pptx'):
                content = read_ppt(file_path)
            else:
                print(f"Unsupported file format: {filename}")
                unread_documents.append(filename)
                continue

            doc = Document(page_content=content, metadata={"filename": filename})
            documents.append(doc)
        except Exception as e:
            print(f"Failed to process {filename}: {e}")
            unread_documents.append(filename)
    return documents, unread_documents

def batch_documents(documents, batch_size=100):
    return [documents[i:i + batch_size] for i in range(0, len(documents), batch_size)]

def embed_documents(document_batches):
    embeddings = []
    embedding_model = VertexAIEmbeddings()
    for batch in document_batches:
        batch_embeddings = embedding_model.embed_documents([doc.page_content for doc in batch])
        embeddings.extend(batch_embeddings)
    return embeddings

def create_and_save_faiss_index(embeddings, documents, index_path='faiss_index.pkl'):
    faiss_index = FAISS()
    faiss_index.add(embeddings, documents)
    with open(index_path, 'wb') as f:
        pickle.dump(faiss_index, f)

def load_faiss_index(index_path='faiss_index.pkl'):
    with open(index_path, 'rb') as f:
        return pickle.load(f)

def perform_rag(query, faiss_index):
    results = faiss_index.search(query)
    return results

# Example usage
directory = 'path_to_your_directory'
documents, unread_documents = load_files(directory)
document_batches = batch_documents(documents)
embeddings = embed_documents(document_batches)
create_and_save_faiss_index(embeddings, documents)

# To retrieve and use the FAISS index
faiss_index = load_faiss_index()
query = "example query text"
results = perform_rag(query, faiss_index)

print("Retrieved Results:")
for result in results:
    print(result.metadata["filename"], result.page_content)
